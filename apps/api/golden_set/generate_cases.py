"""Generate synthetic golden set test cases.

Generates 10 test cases in apps/api/golden_set/cases/
Each case has:
- input.pptx (source content)
- template.pptx (target master/layout)
- notes.md (description)
"""

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

CASES_DIR = Path(__file__).parent / "cases"


def create_simple_deck(path: Path, title: str, content: list[str]):
    """Create a simple deck with Title + Content."""
    prs = Presentation()

    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    if slide.placeholders[1]:
        slide.placeholders[1].text = "Generated Source Deck"

    # Content Slides
    for i, text in enumerate(content):
        layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = f"Slide {i+1}"
        slide.placeholders[1].text = text

    prs.save(str(path))


def create_template_deck(path: Path, theme_color: str = "FF0000"):
    """Create a 'template' deck (standard PPTX but treated as template)."""
    prs = Presentation()
    # Just save a standard blank one, maybe add a shape to distinguish
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "TEMPLATE MASTER"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    textbox.text_frame.text = f"Theme Color: {theme_color}"
    prs.save(str(path))


def generate_cases():
    """Generate 10 test cases."""
    if CASES_DIR.exists():
        shutil.rmtree(CASES_DIR)
    CASES_DIR.mkdir(parents=True)

    cases = [
        ("case_01_basic_text", "Basic Text Transfer", ["Bullet 1", "Bullet 2", "Para 1"]),
        ("case_02_long_content", "Long Content Flow", ["A" * 100, "B" * 200]),
        ("case_03_multi_slide", "Multi Slide Source", ["Slide A", "Slide B", "Slide C", "Slide D"]),
        ("case_04_special_chars", "Special Chars", ["Test @#$%^&*", "Emoji 🚀"]),
        ("case_05_bullets", "Bullet List", ["\n".join([f"- Item {i}" for i in range(5)])]),
        ("case_06_empty_title", "Empty Title", ["Content only", "No title here"]),
        ("case_07_dense_text", "Dense Text Block", ["Start " + "word " * 50 + "End"]),
        ("case_08_short_deck", "One Slide Deck", ["Just one slide"]),
        ("case_09_mixed_content", "Mixed Content", ["Short", "Long " * 20, "Short again"]),
        ("case_10_stress_test", "Stress Test", ["Line\n" * 20]),
    ]

    for name, title, content in cases:
        case_path = CASES_DIR / name
        case_path.mkdir()

        # Input
        create_simple_deck(case_path / "input.pptx", title, content)

        # Template
        create_template_deck(case_path / "template.pptx")

        # Notes
        (case_path / "notes.md").write_text(f"""# {name}

**Category**: {title}
**Expected Outcome**: Content should be preserved exactly.
**Elements**: {len(content) + 1} slides.
""")

        print(f"Generated {name}")


if __name__ == "__main__":
    generate_cases()
