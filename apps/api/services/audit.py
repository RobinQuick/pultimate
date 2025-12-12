from typing import List, Dict, Callable
from ..models import SlideStatus, IssueSeverity, AuditIssue, SlideSummary
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import logging
import os

logger = logging.getLogger(__name__)

class AuditRule:
    def __init__(self, id: str, check_func: Callable, severity: IssueSeverity):
        self.id = id
        self.check_func = check_func
        self.severity = severity

class RuleRegistry:
    def __init__(self):
        self.rules: List[AuditRule] = []

    def register(self, id: str, severity: IssueSeverity):
        def decorator(func):
            self.rules.append(AuditRule(id, func, severity))
            return func
        return decorator

registry = RuleRegistry()

# --- Rules ---

@registry.register("RULE_EMPTY_TITLE", IssueSeverity.WARNING)
def check_empty_title(slide, shape) -> str:
    if shape == slide.shapes.title:
        if not shape.has_text_frame or not shape.text_frame.text.strip():
            return "Title slide has empty title"
    return None

@registry.register("RULE_TEXT_OVERFLOW", IssueSeverity.WARNING)
def check_text_overflow(slide, shape) -> str:
    # Heuristic: if text length > 1000 chars in a placeholder
    if shape.has_text_frame and len(shape.text_frame.text) > 1000:
        return "Text seems too long for placeholder (>1000 chars)"
    return None

@registry.register("RULE_SMARTART_UNSUPPORTED", IssueSeverity.ERROR)
def check_smartart(slide, shape) -> str:
    # python-pptx identifies smartart as GRAPHIC_FRAME broadly, usually inspection of xml required
    # validation simplified: check if it's a graphic frame that is not a table or chart
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
         return "Grouped shapes detected (difficult to preserve exact fidelity)"
    if shape.shape_type == MSO_SHAPE_TYPE.CHART:
        return "Chart detected (requires manual review or rebuild)"
    return None

class AuditEngine:
    def audit_deck(self, pptx_path: str) -> List[SlideSummary]:
        prs = Presentation(pptx_path)
        summaries = []

        for i, slide in enumerate(prs.slides):
            issues = []
            status = SlideStatus.CLEAN

            # Check rules against all shapes
            for shape in slide.shapes:
                for rule in registry.rules:
                    try:
                        msg = rule.check_func(slide, shape)
                        if msg:
                            issues.append(AuditIssue(
                                rule_id=rule.id,
                                severity=rule.severity,
                                message=msg,
                                slide_index=i
                            ))
                            # Update status based on severity logic
                            if rule.severity == IssueSeverity.ERROR:
                                status = SlideStatus.REBUILD
                            elif rule.severity == IssueSeverity.WARNING and status != SlideStatus.REBUILD:
                                status = SlideStatus.REVIEW
                    except Exception as e:
                        logger.error(f"Rule {rule.id} failed on slide {i}: {e}")

            summaries.append(SlideSummary(
                index=i,
                status=status,
                issues=issues
            ))
        
        return summaries

audit_engine = AuditEngine()
