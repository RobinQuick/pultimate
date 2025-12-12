from typing import List, Type, Dict
from pptx import Presentation
from .base import BaseFixer, FixResult
from ...schemas.template_spec import TemplateSpec
from ...services.rules.base import FindingSpec
import logging

logger = logging.getLogger(__name__)

class FixerRegistry:
    _fixers: Dict[str, Type[BaseFixer]] = {}

    @classmethod
    def register(cls, fixer_cls: Type[BaseFixer]):
        cls._fixers[fixer_cls.rule_id] = fixer_cls
        return fixer_cls

    @classmethod
    def get_fixer(cls, rule_id: str) -> Optional[BaseFixer]:
        fixer_cls = cls._fixers.get(rule_id)
        if fixer_cls:
            return fixer_cls()
        return None

class RestyleEngine:
    def apply_fixes(self, pptx_path: str, output_path: str, findings: List[FindingSpec], template: TemplateSpec, mode: str = "SAFE") -> List[FixResult]:
        """
        Loads the PPTX, applies fixes for each finding, saves to output_path.
        """
        prs = Presentation(pptx_path)
        results = []
        
        # Group findings by slide to minimize slide lookups? 
        # Actually findings have slide_index.
        
        # Sort findings by slide index to handle sequential processing if needed
        findings.sort(key=lambda x: x.slide_index)
        
        for finding in findings:
            if finding.slide_index >= len(prs.slides):
                continue
                
            slide = prs.slides[finding.slide_index]
            
            # lookup fixer
            fixer = FixerRegistry.get_fixer(finding.rule_id)
            if not fixer:
                results.append(FixResult(
                    finding_id=finding.rule_id,
                    element_id=finding.element_id or "unknown",
                    action_taken="No fixer registered for this rule",
                    status="SKIPPED",
                    confidence=0.0
                ))
                continue
                
            # Mode Check (Naive implementation: Fixer could check mode, or Registry could filter)
            # For V1: We'll assume all registered fixers are valid for the requested mode 
            # or implemented logic inside Fixer.
            # Example: Layout fixes skipped if SAFE mode.
            if mode == "SAFE" and finding.rule_id == "IMG_QUALITY":
                 # Skip layout changes in SAFE mode?
                 results.append(FixResult(
                    finding_id=finding.rule_id,
                    element_id=finding.element_id or "unknwon",
                    action_taken="Skipped in SAFE mode",
                    status="SKIPPED",
                    confidence=1.0
                ))
                 continue

            try:
                result = fixer.apply(slide, finding, template)
                results.append(result)
            except Exception as e:
                logger.error(f"Fix failed for {finding}: {e}")
                results.append(FixResult(
                    finding_id=finding.rule_id,
                    element_id=finding.element_id or "unknown",
                    action_taken=f"Error: {str(e)}",
                    status="FAILED",
                    confidence=0.0
                ))
        
        prs.save(output_path)
        return results

restyle_engine = RestyleEngine()
