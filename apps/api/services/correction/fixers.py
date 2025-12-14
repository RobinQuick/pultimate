from pptx.dml.color import RGBColor
from pptx.slide import Slide

from schemas.template_spec import TemplateSpec
from services.rules.base import FindingSpec
from .base import BaseFixer, FixResult
from .engine import FixerRegistry


def find_shape_by_id(slide: Slide, shape_id: str):
    # python-pptx shape_id is int, our Spec uses str
    # Iteration needed
    try:
        sid = int(shape_id)
        for shape in slide.shapes:
            if shape.shape_id == sid:
                return shape
    except ValueError:
        pass
    return None

@FixerRegistry.register
class FontFixer(BaseFixer):
    rule_id = "FONT_MISMATCH"

    def apply(self, slide: Slide, finding: FindingSpec, template: TemplateSpec) -> FixResult:
        shape = find_shape_by_id(slide, finding.element_id)
        if not shape or not shape.has_text_frame:
             return FixResult(element_id=finding.element_id, action_taken="Shape not found or has no text", status="FAILED")
        
        # Determine target font
        # Heuristic: If it looks like a Title (based on placeholder or position), use Major. Else Minor.
        # Simple V1: Always Minor (Safety first) or naive check
        target_font = template.theme_fonts.minor
        
        # Apply to all runs
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                run.font.name = target_font
        
        return FixResult(
            element_id=finding.element_id,
            action_taken=f"Changed font to {target_font}",
            status="SUCCESS"
        )

@FixerRegistry.register
class ColorFixer(BaseFixer):
    rule_id = "COLOR_MISMATCH"

    def apply(self, slide: Slide, finding: FindingSpec, template: TemplateSpec) -> FixResult:
        shape = find_shape_by_id(slide, finding.element_id)
        if not shape or not shape.has_text_frame:
             return FixResult(element_id=finding.element_id, action_taken="Shape not found", status="FAILED")
        
        # Target color: Dark1 (Main text color usually)
        # In real V2: Find nearest semantic color
        target_color = template.theme_colors.dark1
        if not target_color:
            # Fallback black
            r, g, b = 0, 0, 0
        else:
            r, g, b = target_color.r, target_color.g, target_color.b
            
        for p in shape.text_frame.paragraphs:
            for run in p.runs:
                run.font.color.rgb = RGBColor(r, g, b)
                
        return FixResult(
            element_id=finding.element_id,
            action_taken=f"Changed color to Theme Dark1 ({r},{g},{b})",
            status="SUCCESS"
        )

@FixerRegistry.register
class ImageFixer(BaseFixer):
    rule_id = "IMG_QUALITY" # Stretched

    def apply(self, slide: Slide, finding: FindingSpec, template: TemplateSpec) -> FixResult:
        shape = find_shape_by_id(slide, finding.element_id)
        if not shape: # or shape type not picture
             return FixResult(element_id=finding.element_id, action_taken="Shape not found", status="FAILED")

        # Reset aspect ratio logic stub
        # python-pptx doesn't easily expose "original size" to calculate pure resets without complexity
        # But we can try setting scale_y = scale_x if we assume one dimension is "correct"
        # Or usually it's better to NOT touch in V1 Safe Mode if we aren't sure.
        # Let's assume we just want to ensure it's not distorted.
        
        # STUB: Just logging action without destructive change for V1 safety unless we have data
        return FixResult(
            element_id=finding.element_id,
            action_taken="Marked for manual review (Aspect Ratio reset complex in V1)",
            status="SKIPPED"
        )
