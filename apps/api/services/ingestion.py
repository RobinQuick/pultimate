from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.shapes.placeholder import BasePlaceholder
from ..schemas.template_spec import TemplateSpec,  MasterSpec, LayoutSpec, PlaceholderSpec, ThemeColors, ThemeFonts, RgbColor
import logging

logger = logging.getLogger(__name__)

class TemplateIngestor:
    def ingest(self, pptx_path: str) -> TemplateSpec:
        prs = Presentation(pptx_path)
        
        # 1. Extract Theme (Simplified approximation as python-pptx extraction of theme XML is non-trivial high-level)
        # Real implementation would parse ppt/theme/theme1.xml
        # Here we stub/approximate or use what we can if available.
        # python-pptx doesn't expose clean theme API yet.
        theme_colors = ThemeColors() 
        theme_fonts = ThemeFonts(major="Arial", minor="Arial") # Default stub
        
        masters = []
        
        # 2. Iter Masters
        for i, master in enumerate(prs.slide_masters):
            layout_specs = []
            
            # 3. Iter Layouts
            for j, layout in enumerate(master.slide_layouts):
                placeholders = []
                for shape in layout.placeholders:
                    ph_spec = PlaceholderSpec(
                        idx=shape.placeholder_format.idx,
                        type=str(shape.placeholder_format.type), # Enum to str
                        name=shape.name,
                        left=shape.left,
                        top=shape.top,
                        width=shape.width,
                        height=shape.height
                    )
                    placeholders.append(ph_spec)
                
                layout_specs.append(LayoutSpec(
                    index=j,
                    name=layout.name,
                    placeholders=placeholders
                ))
                
            masters.append(MasterSpec(
                id=i,
                name=master.name or f"Master {i}",
                layouts=layout_specs
            ))
            
        return TemplateSpec(
            name="Ingested Template",
            theme_colors=theme_colors,
            theme_fonts=theme_fonts,
            masters=masters
        )

ingestor = TemplateIngestor()
