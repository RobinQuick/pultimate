import logging

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from ..schemas.slide_spec import BoundingBox, DeckSpec, ElementSpec, SlideSpec, SlideStats, TextStyle

logger = logging.getLogger(__name__)

EMU_PER_PX = 9525 # Approx 96 DPI

def emu_to_px(emu_val) -> float:
    if emu_val is None: return 0.0
    return round(float(emu_val) / EMU_PER_PX, 2)

class DeckParser:
    def parse(self, pptx_path: str, filename: str) -> DeckSpec:
        prs = Presentation(pptx_path)
        
        slides = []
        for i, slide in enumerate(prs.slides):
            slide_spec = self._parse_slide(slide, i)
            slides.append(slide_spec)
            
        return DeckSpec(
            filename=filename,
            slide_count=len(slides),
            slides=slides
        )

    def _parse_slide(self, slide, index: int) -> SlideSpec:
        elements = []
        used_fonts = set()
        used_colors = set()
        suspects = []
        
        # Iterate shapes
        # Note: z-order is implicit in python-pptx by iteration order (back to front)
        for z, shape in enumerate(slide.shapes):
            element = self._parse_shape(shape, z)
            if element:
                elements.append(element)
                
                # Gather Stats
                if element.text_style and element.text_style.font_family:
                    used_fonts.add(element.text_style.font_family)
                if element.text_style and element.text_style.color_hex:
                    used_colors.add(element.text_style.color_hex)
                if element.is_stretched:
                    suspects.append(f"Stretched Image: {element.name}")

        return SlideSpec(
            index=index,
            layout_name=slide.slide_layout.name,
            elements=elements,
            stats=SlideStats(
                used_fonts=list(used_fonts),
                used_colors=list(used_colors)
            ),
            suspect_issues=suspects
        )

    def _parse_shape(self, shape, z_order: int) -> Optional[ElementSpec]:
        # Handle Groups (Recursion could be added here, currently flattened or skipped)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            # For V1 simplicity, treating group as bounding box container
            pass

        bbox = BoundingBox(
            x=emu_to_px(shape.left),
            y=emu_to_px(shape.top),
            width=emu_to_px(shape.width),
            height=emu_to_px(shape.height)
        )
        
        elem = ElementSpec(
            id=str(shape.shape_id),
            type=str(shape.shape_type),
            name=shape.name,
            bbox=bbox,
            rotation=shape.rotation,
            z_order=z_order
        )
        
        # Text extraction
        if hasattr(shape, "text") and shape.text.strip():
            elem.text_content = shape.text
            # Extract first paragraph style as representative
            if hasattr(shape, "text_frame") and shape.text_frame.paragraphs:
                p = shape.text_frame.paragraphs[0]
                run = p.runs[0] if p.runs else None
                
                font_family = "Unknown"
                if run and run.font.name:
                    font_family = run.font.name
                elif hasattr(p, "font") and p.font.name:
                    font_family = p.font.name
                
                # Iterate runs to verify consistency? Keep simple for V1
                elem.text_style = TextStyle(
                    font_family=font_family,
                    font_size=emu_to_px(run.font.size) if run and run.font.size else 0, # Actually font size in pptx is confusing (Pt vs Emu). python-pptx creates Pt object.
                    is_bold=run.font.bold if run else False,
                    is_italic=run.font.italic if run else False
                )
                
                # Fix font size conversion if it's Pt
                if run and run.font.size:
                     # run.font.size is Length object (EMU)
                     # 1 Pt = 12700 EMU
                     elem.text_style.font_size = round(run.font.size.pt, 1)

        # Image extraction stats
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            elem.type = "IMAGE"
            # Stretched check
            # python-pptx doesn't give easily original image dimensions without opening the blob
            # Stub logic: assume valid for now or inspect image blob in V2
            pass

        return elem

deck_parser = DeckParser()
