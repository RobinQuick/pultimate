"""Rebuild service - Core logic for deck rebuilding.

NO-GEN POLICY: This service copies content, never generates.
"""

import logging
import tempfile
import uuid
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER
from pptx.util import Emu

from schemas.mapping_schema import (
    BoundingBox,
    DeckElement,
    DeckElementList,
    ElementType,
    MappingAction,
    MappingResult,
    PlaceholderType,
    TemplatePlaceholder,
    TemplatePlaceholderList,
)

logger = logging.getLogger(__name__)


# =============================================================================
# PARSE DECK ELEMENTS
# =============================================================================


def _get_element_type(shape) -> ElementType:
    """Determine element type from shape."""
    if shape.has_table:
        return ElementType.TABLE
    if shape.has_chart:
        return ElementType.CHART
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return ElementType.IMAGE

    # Check for placeholder type
    if shape.is_placeholder:
        ph_type = shape.placeholder_format.type if shape.placeholder_format else None
        if ph_type in (PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.VERTICAL_TITLE):
            return ElementType.TITLE
        if ph_type in (PP_PLACEHOLDER.BODY, PP_PLACEHOLDER.OBJECT, PP_PLACEHOLDER.VERTICAL_BODY):
            return ElementType.BODY

    if hasattr(shape, "text_frame") and shape.has_text_frame:
        return ElementType.BODY

    return ElementType.OTHER


def _extract_text_preview(shape, max_len: int = 200) -> str | None:
    """Extract first N chars of text from shape."""
    if not hasattr(shape, "text_frame") or not shape.has_text_frame:
        return None
    try:
        text = shape.text_frame.text.strip()
        return text[:max_len] if text else None
    except Exception:
        return None


def _emu_to_pt(emu: int | None) -> float:
    """Convert EMU to points."""
    if emu is None:
        return 0.0
    return Emu(emu).pt


def parse_deck_elements(pptx_path: str | Path) -> DeckElementList:
    """Parse a deck and extract all elements with stable IDs.

    Each element gets a stable ID based on:
    - slide index
    - shape ID within slide

    Returns:
        DeckElementList with all elements and slide count
    """
    prs = Presentation(str(pptx_path))
    elements = []

    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # Generate stable element ID
            element_id = f"slide_{slide_idx}_shape_{shape.shape_id}"

            # Get bounding box
            bbox = BoundingBox(
                x=_emu_to_pt(shape.left),
                y=_emu_to_pt(shape.top),
                width=_emu_to_pt(shape.width),
                height=_emu_to_pt(shape.height),
            )

            element = DeckElement(
                element_id=element_id,
                slide_index=slide_idx,
                element_type=_get_element_type(shape),
                name=shape.name,
                bbox=bbox,
                text_preview=_extract_text_preview(shape),
                has_image=shape.shape_type == MSO_SHAPE_TYPE.PICTURE,
                has_table=shape.has_table,
                has_chart=shape.has_chart,
            )
            elements.append(element)

    logger.info(f"Parsed {len(elements)} elements from {len(prs.slides)} slides")

    return DeckElementList(elements=elements, slide_count=len(prs.slides))


# =============================================================================
# PARSE TEMPLATE PLACEHOLDERS
# =============================================================================


def _get_placeholder_type(ph) -> PlaceholderType:
    """Map python-pptx placeholder type to our enum."""
    try:
        ph_type = ph.placeholder_format.type
    except Exception:
        return PlaceholderType.OTHER

    mapping = {
        PP_PLACEHOLDER.TITLE: PlaceholderType.TITLE,
        PP_PLACEHOLDER.CENTER_TITLE: PlaceholderType.TITLE,
        PP_PLACEHOLDER.VERTICAL_TITLE: PlaceholderType.TITLE,
        PP_PLACEHOLDER.SUBTITLE: PlaceholderType.SUBTITLE,
        PP_PLACEHOLDER.BODY: PlaceholderType.BODY,
        PP_PLACEHOLDER.OBJECT: PlaceholderType.CONTENT,
        PP_PLACEHOLDER.CHART: PlaceholderType.CHART,
        PP_PLACEHOLDER.TABLE: PlaceholderType.TABLE,
        PP_PLACEHOLDER.PICTURE: PlaceholderType.PICTURE,
        PP_PLACEHOLDER.FOOTER: PlaceholderType.FOOTER,
        PP_PLACEHOLDER.SLIDE_NUMBER: PlaceholderType.SLIDE_NUMBER,
        PP_PLACEHOLDER.DATE: PlaceholderType.DATE,
    }

    return mapping.get(ph_type, PlaceholderType.OTHER)


def parse_template_placeholders(template_path: str | Path) -> TemplatePlaceholderList:
    """Parse a template and extract all placeholders.

    Returns:
        TemplatePlaceholderList with all placeholders and layout count
    """
    prs = Presentation(str(template_path))
    placeholders = []

    for master in prs.slide_masters:
        for layout_idx, layout in enumerate(master.slide_layouts):
            layout_name = layout.name or f"Layout {layout_idx}"

            for shape in layout.placeholders:
                # Generate stable placeholder ID
                ph_idx = shape.placeholder_format.idx if shape.placeholder_format else None
                placeholder_id = f"layout_{layout_idx}_ph_{ph_idx or shape.shape_id}"

                bbox = BoundingBox(
                    x=_emu_to_pt(shape.left),
                    y=_emu_to_pt(shape.top),
                    width=_emu_to_pt(shape.width),
                    height=_emu_to_pt(shape.height),
                )

                placeholder = TemplatePlaceholder(
                    placeholder_id=placeholder_id,
                    layout_name=layout_name,
                    layout_index=layout_idx,
                    placeholder_type=_get_placeholder_type(shape),
                    bbox=bbox,
                    idx=ph_idx,
                )
                placeholders.append(placeholder)

    layout_count = sum(len(m.slide_layouts) for m in prs.slide_masters)
    logger.info(f"Parsed {len(placeholders)} placeholders from {layout_count} layouts")

    return TemplatePlaceholderList(placeholders=placeholders, layout_count=layout_count)


# =============================================================================
# APPLY MAPPING (NO AI - PURE COPY)
# =============================================================================


class RebuildResult:
    """Result of rebuild operation."""

    def __init__(self):
        self.output_path: Path | None = None
        self.slides_created: int = 0
        self.elements_mapped: int = 0
        self.elements_skipped: int = 0
        self.warnings: list[str] = []
        self.errors: list[str] = []


def apply_mapping(
    deck_path: str | Path,
    template_path: str | Path,
    mapping: MappingResult,
    output_dir: str | Path | None = None,
) -> RebuildResult:
    """Apply mapping to rebuild deck using template.

    NO-GEN POLICY: Content is ONLY copied, never generated.

    Args:
        deck_path: Path to source deck
        template_path: Path to template
        mapping: Validated mapping result
        output_dir: Output directory (temp if None)

    Returns:
        RebuildResult with output path and stats
    """
    result = RebuildResult()

    # Create output directory
    if output_dir is None:
        output_dir = Path(tempfile.mkdtemp(prefix="rebuild_"))
    else:
        output_dir = Path(output_dir)
        output_dir.mkdir(parents=True, exist_ok=True)

    output_path = output_dir / f"rebuilt_{uuid.uuid4().hex[:8]}.pptx"

    try:
        # Load source deck and template
        source_prs = Presentation(str(deck_path))
        template_prs = Presentation(str(template_path))

        # Get all layouts from template
        layouts = []
        for master in template_prs.slide_masters:
            layouts.extend(master.slide_layouts)

        if not layouts:
            result.errors.append("Template has no layouts")
            return result

        # Build shape lookup from source deck
        # shape_lookup[slide_idx][shape_id] = shape
        shape_lookup: dict[int, dict[int, any]] = {}
        for slide_idx, slide in enumerate(source_prs.slides):
            shape_lookup[slide_idx] = {}
            for shape in slide.shapes:
                shape_lookup[slide_idx][shape.shape_id] = shape

        # Process each output slide
        for slide_mapping in mapping.slide_mappings:
            layout_idx = min(slide_mapping.layout_index, len(layouts) - 1)
            layout = layouts[layout_idx]

            # Add new slide with layout
            new_slide = template_prs.slides.add_slide(layout)
            result.slides_created += 1

            # Build placeholder lookup for this slide
            ph_lookup = {}
            for shape in new_slide.placeholders:
                ph_idx = shape.placeholder_format.idx if shape.placeholder_format else None
                ph_id = f"layout_{layout_idx}_ph_{ph_idx or shape.shape_id}"
                ph_lookup[ph_id] = shape

            # Apply each element mapping
            for em in slide_mapping.element_mappings:
                if em.action == MappingAction.SKIP:
                    result.elements_skipped += 1
                    continue

                if em.action == MappingAction.MAP:
                    # Parse source element ID
                    try:
                        parts = em.source_element_id.split("_")
                        src_slide_idx = int(parts[1])
                        src_shape_id = int(parts[3])
                    except (IndexError, ValueError):
                        result.warnings.append(f"Invalid element ID: {em.source_element_id}")
                        continue

                    # Get source shape
                    source_shape = shape_lookup.get(src_slide_idx, {}).get(src_shape_id)
                    if not source_shape:
                        result.warnings.append(f"Source shape not found: {em.source_element_id}")
                        continue

                    # Get target placeholder
                    target_ph = ph_lookup.get(em.target_placeholder_id)
                    if not target_ph:
                        result.warnings.append(f"Target placeholder not found: {em.target_placeholder_id}")
                        continue

                    # Copy content (NO GENERATION)
                    _copy_shape_content(source_shape, target_ph)
                    result.elements_mapped += 1

        # Save output
        template_prs.save(str(output_path))
        result.output_path = output_path

        logger.info(
            f"Rebuild complete: {result.slides_created} slides, "
            f"{result.elements_mapped} mapped, {result.elements_skipped} skipped"
        )

    except Exception as e:
        result.errors.append(str(e))
        logger.error(f"Rebuild failed: {e}")

    return result


def _copy_shape_content(source_shape, target_shape) -> None:
    """Copy content from source shape to target shape.

    NO-GEN POLICY: Only copies existing content, never generates.
    """
    try:
        # Copy text content
        if hasattr(source_shape, "text_frame") and source_shape.has_text_frame:
            if hasattr(target_shape, "text_frame"):
                # Clear target
                target_tf = target_shape.text_frame
                for para in list(target_tf.paragraphs)[1:]:
                    p = para._p
                    p.getparent().remove(p)

                # Copy paragraphs
                for para_idx, src_para in enumerate(source_shape.text_frame.paragraphs):
                    if para_idx >= len(target_tf.paragraphs):
                        target_tf.paragraphs[0].add_run()  # Simplified

                    target_para = target_tf.paragraphs[min(para_idx, len(target_tf.paragraphs) - 1)]
                    target_para.text = src_para.text

                    # Copy font properties if runs exist
                    if src_para.runs and target_para.runs:
                        for run_idx, src_run in enumerate(src_para.runs):
                            if run_idx < len(target_para.runs):
                                target_run = target_para.runs[run_idx]
                                if src_run.font.name:
                                    target_run.font.name = src_run.font.name
                                if src_run.font.size:
                                    target_run.font.size = src_run.font.size
                                if src_run.font.bold is not None:
                                    target_run.font.bold = src_run.font.bold

    except Exception as e:
        logger.warning(f"Content copy failed: {e}")


# =============================================================================
# DOWNLOAD / UPLOAD HELPERS
# =============================================================================


async def download_from_r2(s3_key: str, local_path: Path) -> bool:
    """Download file from R2 to local path."""
    from services.storage import storage

    try:
        await storage.download_file(s3_key, str(local_path))
        return True
    except Exception as e:
        logger.error(f"Failed to download {s3_key}: {e}")
        return False


async def upload_to_r2(local_path: Path, s3_key: str) -> bool:
    """Upload file from local path to R2."""
    from services.storage import storage

    try:
        # Read file and upload
        with open(local_path, "rb") as f:
            await storage.upload_bytes(f.read(), s3_key)
        return True
    except Exception as e:
        logger.error(f"Failed to upload to {s3_key}: {e}")
        return False
