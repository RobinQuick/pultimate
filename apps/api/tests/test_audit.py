import pytest
import os
from pptx import Presentation
from apps.api.services.audit import audit_engine, SlideStatus, IssueSeverity

def create_synthetic_pptx(path, title_text=None, body_text=None):
    prs = Presentation()
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    
    if title_text is not None:
        slide.shapes.title.text = title_text
    
    if body_text is not None:
        # Add a text box or use subtitle
        pass # simplified for this test
        
    prs.save(path)

def test_audit_empty_title(tmp_path):
    pptx_path = tmp_path / "test.pptx"
    create_synthetic_pptx(pptx_path, title_text="") # Empty title
    
    summaries = audit_engine.audit_deck(str(pptx_path))
    
    assert len(summaries) == 1
    assert summaries[0].status == SlideStatus.REVIEW # Warnings mapped to REVIEW in code currently?
    # Wait, my code said: ERROR->REBUILD, WARNING->REVIEW.
    # TITLE_EMPTY is WARNING. So yes, REVIEW.
    
    issues = summaries[0].issues
    assert len(issues) > 0
    assert issues[0].rule_id == "RULE_EMPTY_TITLE"

def test_audit_clean_title(tmp_path):
    pptx_path = tmp_path / "clean.pptx"
    create_synthetic_pptx(pptx_path, title_text="Valid Title")
    
    summaries = audit_engine.audit_deck(str(pptx_path))
    
    assert len(summaries) == 1
    assert summaries[0].status == SlideStatus.CLEAN
    assert len(summaries[0].issues) == 0
