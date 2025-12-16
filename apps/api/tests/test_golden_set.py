"""Pytest tests for golden set validation.

Tests:
- test_content_preserved: exact text match checks
- test_layout_sanity: no shapes outside slide bounds
- test_no_gen_policy: LLM output has no content generation fields
"""

import asyncio
from pathlib import Path

from pptx import Presentation

from golden_set.runner import (
    check_layout_sanity,
    check_no_gen_policy,
    compute_text_hash,
    extract_text_from_pptx,
)
from schemas.mapping_schema import (
    BoundingBox,
    DeckElement,
    ElementMapping,
    ElementType,
    LLMConfig,
    MappingAction,
    MappingResult,
    PlaceholderType,
    SlideMapping,
    TemplatePlaceholder,
)


# =============================================================================
# CONTENT PRESERVATION TESTS
# =============================================================================


def create_test_pptx(path: Path, texts: list[str]) -> None:
    """Create a test PPTX with given texts."""
    prs = Presentation()
    for _i, text in enumerate(texts):
        slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
        left = top = 100000
        width = height = 500000
        shape = slide.shapes.add_textbox(left, top, width, height)
        shape.text_frame.text = text
    prs.save(str(path))


def test_extract_text_from_pptx(tmp_path):
    """Text extraction works correctly."""
    pptx_path = tmp_path / "test.pptx"
    expected_texts = ["Hello World", "Test Content", "More Text"]
    create_test_pptx(pptx_path, expected_texts)

    extracted = extract_text_from_pptx(pptx_path)

    for text in expected_texts:
        assert any(text in e for e in extracted), f"'{text}' not found in extracted texts"


def test_content_hash_consistency(tmp_path):
    """Same content produces same hash."""
    pptx1 = tmp_path / "test1.pptx"
    pptx2 = tmp_path / "test2.pptx"

    texts = ["Content A", "Content B"]
    create_test_pptx(pptx1, texts)
    create_test_pptx(pptx2, texts)

    hash1 = compute_text_hash(extract_text_from_pptx(pptx1))
    hash2 = compute_text_hash(extract_text_from_pptx(pptx2))

    assert hash1 == hash2


def test_content_hash_detects_changes(tmp_path):
    """Different content produces different hash."""
    pptx1 = tmp_path / "test1.pptx"
    pptx2 = tmp_path / "test2.pptx"

    create_test_pptx(pptx1, ["Original Text"])
    create_test_pptx(pptx2, ["Modified Text"])

    hash1 = compute_text_hash(extract_text_from_pptx(pptx1))
    hash2 = compute_text_hash(extract_text_from_pptx(pptx2))

    assert hash1 != hash2


# =============================================================================
# LAYOUT SANITY TESTS
# =============================================================================


def test_layout_sanity_valid_pptx(tmp_path):
    """Valid PPTX passes layout sanity check."""
    pptx_path = tmp_path / "valid.pptx"
    create_test_pptx(pptx_path, ["Normal content"])

    is_valid, issues = check_layout_sanity(pptx_path)

    assert is_valid is True
    assert len(issues) == 0


def test_layout_sanity_detects_negative_position(tmp_path):
    """Layout check detects shapes with negative positions."""
    pptx_path = tmp_path / "invalid.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Add shape with negative position
    from pptx.util import Emu

    shape = slide.shapes.add_textbox(Emu(-100000), Emu(-100000), Emu(100000), Emu(100000))
    shape.text_frame.text = "Invalid position"

    prs.save(str(pptx_path))

    is_valid, issues = check_layout_sanity(pptx_path)

    assert is_valid is False
    assert any("negative position" in issue for issue in issues)


# =============================================================================
# NO-GEN POLICY TESTS
# =============================================================================


def test_no_gen_policy_valid_mapping():
    """Valid mapping passes NO-GEN check."""
    mapping = {
        "slide_mappings": [
            {
                "output_slide_index": 0,
                "layout_index": 0,
                "layout_name": "Title Slide",
                "element_mappings": [
                    {
                        "source_element_id": "slide_0_shape_2",
                        "target_placeholder_id": "layout_0_ph_0",
                        "action": "MAP",
                    }
                ],
            }
        ],
        "skipped_elements": [],
        "warnings": [],
    }

    is_valid, violations = check_no_gen_policy(mapping)

    assert is_valid is True
    assert len(violations) == 0


def test_no_gen_policy_detects_forbidden_fields():
    """NO-GEN check detects forbidden content generation fields."""
    invalid_mapping = {
        "slide_mappings": [
            {
                "output_slide_index": 0,
                "element_mappings": [
                    {
                        "source_element_id": "slide_0_shape_2",
                        "generated_text": "This should not exist!",  # FORBIDDEN
                    }
                ],
            }
        ],
    }

    is_valid, violations = check_no_gen_policy(invalid_mapping)

    assert is_valid is False
    assert any("generated_text" in v for v in violations)


def test_no_gen_policy_nested_forbidden_field():
    """NO-GEN check detects deeply nested forbidden fields."""
    invalid_mapping = {
        "data": {
            "nested": {
                "deeply": {
                    "new_content": "Generated content here",  # FORBIDDEN deep in structure
                }
            }
        }
    }

    is_valid, violations = check_no_gen_policy(invalid_mapping)

    assert is_valid is False
    assert len(violations) > 0


# =============================================================================
# MAPPING SCHEMA NO-GEN TESTS
# =============================================================================


def test_mapping_schema_has_no_content_fields():
    """Verify ElementMapping schema has no content generation fields."""
    mapping = ElementMapping(
        source_element_id="test",
        target_placeholder_id="placeholder",
        action=MappingAction.MAP,
    )

    # Verify no content fields exist
    fields = set(mapping.model_fields.keys())
    forbidden = {"generated_text", "new_content", "modified_text", "ai_content"}

    assert fields.isdisjoint(forbidden), f"Found forbidden fields: {fields & forbidden}"


def test_mapping_result_only_contains_ids():
    """MappingResult only contains IDs and mapping decisions."""
    result = MappingResult(
        slide_mappings=[
            SlideMapping(
                output_slide_index=0,
                layout_index=0,
                layout_name="Test Layout",
                element_mappings=[
                    ElementMapping(
                        source_element_id="slide_0_shape_1",
                        target_placeholder_id="layout_0_ph_0",
                        action=MappingAction.MAP,
                    )
                ],
            )
        ],
        skipped_elements=["slide_0_shape_2"],
        warnings=["Test warning"],
    )

    # Dump to dict and verify no content
    data = result.model_dump()

    def find_long_strings(d, path="", max_len=100):
        """Find any suspiciously long strings that might be content."""
        long_strings = []
        if isinstance(d, dict):
            for k, v in d.items():
                find_long_strings(v, f"{path}.{k}", max_len)
        elif isinstance(d, list):
            for i, item in enumerate(d):
                find_long_strings(item, f"{path}[{i}]", max_len)
        elif isinstance(d, str) and len(d) > max_len:
            long_strings.append((path, d))
        return long_strings

    long_content = find_long_strings(data)
    assert len(long_content) == 0, f"Found long strings (potential content): {long_content}"


# =============================================================================
# MOCK LLM INTEGRATION TEST
# =============================================================================


def test_mock_llm_returns_valid_mapping():
    """Mock LLM returns schema-valid mapping with no content."""
    from services.llm_service import call_llm_for_mapping

    elements = [
        DeckElement(
            element_id="slide_0_shape_1",
            slide_index=0,
            element_type=ElementType.TITLE,
            bbox=BoundingBox(x=0, y=0, width=100, height=50),
        )
    ]

    placeholders = [
        TemplatePlaceholder(
            placeholder_id="layout_0_ph_0",
            layout_name="Title",
            layout_index=0,
            placeholder_type=PlaceholderType.TITLE,
            bbox=BoundingBox(x=0, y=0, width=100, height=50),
        )
    ]

    config = LLMConfig(provider="mock", model="mock")
    result = call_llm_for_mapping(elements, placeholders, config)

    # Verify it's a valid mapping
    assert isinstance(result, MappingResult)

    # Verify NO-GEN policy
    is_valid, violations = check_no_gen_policy(result.model_dump())
    assert is_valid, f"NO-GEN violations: {violations}"
