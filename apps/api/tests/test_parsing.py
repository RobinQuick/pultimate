from apps.api.services.parser import deck_parser
from pptx import Presentation
from pptx.util import Pt


def create_synthetic_deck(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0]) # Title slide
    
    # Add Title
    title = slide.shapes.title
    title.text = "Test Parser"
    # Basic styling
    p = title.text_frame.paragraphs[0]
    p.font.name = "Arial"
    p.font.size = Pt(24)
    p.font.bold = True
    
    # Add a Box
    left = top = width = height = Pt(100)
    shape = slide.shapes.add_shape(
        1, # MSO_SHAPE.RECTANGLE (approx raw int)
        left, top, width, height
    )
    shape.name = "MyRect"
    
    prs.save(path)

def test_deck_parsing(tmp_path):
    pptx_path = tmp_path / "test_deck.pptx"
    create_synthetic_deck(pptx_path)
    
    spec = deck_parser.parse(str(pptx_path), "test_deck.pptx")
    
    assert spec.filename == "test_deck.pptx"
    assert spec.slide_count == 1
    
    slide = spec.slides[0]
    # Check Title
    # Title is usually shape 0 but placeholders order varies.
    # We find by text
    title_elem = next(e for e in slide.elements if e.text_content == "Test Parser")
    assert title_elem
    assert title_elem.text_style.font_family == "Arial"
    assert title_elem.text_style.font_size == 24.0
    assert title_elem.text_style.is_bold is True
    
    # Check Stats
    assert "Arial" in slide.stats.used_fonts
