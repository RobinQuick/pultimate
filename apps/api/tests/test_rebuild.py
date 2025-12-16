"""Tests for rebuild engine (Phase B).

NO-GEN POLICY: Tests verify content is copied, never generated.
"""

import pytest
from pptx import Presentation

from schemas.mapping_schema import (
    BoundingBox,
    DeckElement,
    ElementMapping,
    ElementType,
    LLMConfig,
    MappingAction,
    MappingResult,
    PlaceholderType,
    SlideMapping,
    TemplatePlaceholder,
    validate_mapping_against_inputs,
)


# =============================================================================
# MAPPING SCHEMA TESTS
# =============================================================================


def test_element_mapping_requires_placeholder_for_map():
    """MAP action must have target_placeholder_id."""
    with pytest.raises(ValueError, match="target_placeholder_id required"):
        ElementMapping(
            source_element_id="slide_0_shape_1",
            target_placeholder_id=None,  # Invalid for MAP
            action=MappingAction.MAP,
        )


def test_element_mapping_allows_none_for_skip():
    """SKIP action can have null placeholder."""
    mapping = ElementMapping(
        source_element_id="slide_0_shape_1",
        target_placeholder_id=None,
        action=MappingAction.SKIP,
    )
    assert mapping.action == MappingAction.SKIP


def test_mapping_result_requires_slides():
    """MappingResult must have at least one slide mapping."""
    with pytest.raises(ValueError, match="cannot be empty"):
        MappingResult(slide_mappings=[], skipped_elements=[], warnings=[])


def test_valid_mapping_result():
    """Valid MappingResult passes validation."""
    result = MappingResult(
        slide_mappings=[
            SlideMapping(
                output_slide_index=0,
                layout_index=0,
                layout_name="Title Slide",
                element_mappings=[
                    ElementMapping(
                        source_element_id="slide_0_shape_2",
                        target_placeholder_id="layout_0_ph_0",
                        action=MappingAction.MAP,
                        target_layout_index=0,
                        target_slide_index=0,
                    )
                ],
            )
        ],
        skipped_elements=[],
        warnings=[],
    )

    assert len(result.slide_mappings) == 1
    assert len(result.slide_mappings[0].element_mappings) == 1


def test_validate_mapping_against_inputs_valid():
    """Validation passes for correct references."""
    elements = [
        DeckElement(
            element_id="slide_0_shape_2",
            slide_index=0,
            element_type=ElementType.TITLE,
            bbox=BoundingBox(x=0, y=0, width=100, height=50),
        )
    ]

    placeholders = [
        TemplatePlaceholder(
            placeholder_id="layout_0_ph_0",
            layout_name="Title Slide",
            layout_index=0,
            placeholder_type=PlaceholderType.TITLE,
            bbox=BoundingBox(x=0, y=0, width=100, height=50),
        )
    ]

    mapping = MappingResult(
        slide_mappings=[
            SlideMapping(
                output_slide_index=0,
                layout_index=0,
                layout_name="Title Slide",
                element_mappings=[
                    ElementMapping(
                        source_element_id="slide_0_shape_2",
                        target_placeholder_id="layout_0_ph_0",
                        action=MappingAction.MAP,
                    )
                ],
            )
        ],
    )

    errors = validate_mapping_against_inputs(mapping, elements, placeholders)
    assert errors == []


def test_validate_mapping_against_inputs_invalid_element():
    """Validation fails for unknown element ID."""
    elements = [
        DeckElement(
            element_id="slide_0_shape_2",
            slide_index=0,
            element_type=ElementType.TITLE,
            bbox=BoundingBox(x=0, y=0, width=100, height=50),
        )
    ]

    placeholders = [
        TemplatePlaceholder(
            placeholder_id="layout_0_ph_0",
            layout_name="Title Slide",
            layout_index=0,
            placeholder_type=PlaceholderType.TITLE,
            bbox=BoundingBox(x=0, y=0, width=100, height=50),
        )
    ]

    mapping = MappingResult(
        slide_mappings=[
            SlideMapping(
                output_slide_index=0,
                layout_index=0,
                layout_name="Title Slide",
                element_mappings=[
                    ElementMapping(
                        source_element_id="slide_0_shape_999",  # Invalid
                        target_placeholder_id="layout_0_ph_0",
                        action=MappingAction.MAP,
                    )
                ],
            )
        ],
    )

    errors = validate_mapping_against_inputs(mapping, elements, placeholders)
    assert len(errors) == 1
    assert "Unknown source element" in errors[0]


# =============================================================================
# LLM SERVICE TESTS (WITH MOCK)
# =============================================================================


def test_llm_mock_provider():
    """Mock LLM provider returns valid mapping."""
    from services.llm_service import call_llm_for_mapping

    elements = [
        DeckElement(
            element_id="slide_0_shape_2",
            slide_index=0,
            element_type=ElementType.TITLE,
            bbox=BoundingBox(x=0, y=0, width=100, height=50),
            text_preview="Test Title",
        )
    ]

    placeholders = [
        TemplatePlaceholder(
            placeholder_id="layout_0_ph_0",
            layout_name="Title Slide",
            layout_index=0,
            placeholder_type=PlaceholderType.TITLE,
            bbox=BoundingBox(x=0, y=0, width=100, height=50),
        )
    ]

    config = LLMConfig(provider="mock", model="mock-model")
    result = call_llm_for_mapping(elements, placeholders, config)

    assert isinstance(result, MappingResult)
    assert len(result.slide_mappings) >= 1
    # Mock always adds a warning
    assert any("Mock LLM" in w for w in result.warnings)


# =============================================================================
# PARSE SERVICE TESTS
# =============================================================================


def test_parse_deck_elements(tmp_path):
    """parse_deck_elements extracts elements."""
    from services.rebuild_service import parse_deck_elements

    # Create test pptx
    pptx_path = tmp_path / "test.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Test Title"
    prs.save(str(pptx_path))

    result = parse_deck_elements(pptx_path)

    assert result.slide_count == 1
    assert len(result.elements) >= 1

    # Find title element
    title_elements = [e for e in result.elements if e.element_type == ElementType.TITLE]
    assert len(title_elements) >= 1


def test_parse_template_placeholders(tmp_path):
    """parse_template_placeholders extracts placeholders."""
    from services.rebuild_service import parse_template_placeholders

    # Create test template (just a pptx for simplicity)
    template_path = tmp_path / "template.pptx"
    prs = Presentation()
    prs.save(str(template_path))

    result = parse_template_placeholders(template_path)

    assert result.layout_count >= 1
    # Should have some placeholders from default layouts


# =============================================================================
# NO-GEN POLICY TESTS
# =============================================================================


def test_no_gen_policy_no_content_fields():
    """ElementMapping has no fields for generated content."""
    # This test verifies the schema design enforces NO-GEN policy
    mapping = ElementMapping(
        source_element_id="slide_0_shape_2",
        target_placeholder_id="layout_0_ph_0",
        action=MappingAction.MAP,
    )

    # Should not have any content fields
    assert not hasattr(mapping, "generated_text")
    assert not hasattr(mapping, "new_content")
    assert not hasattr(mapping, "modified_text")

    # Only has mapping fields
    assert hasattr(mapping, "source_element_id")
    assert hasattr(mapping, "target_placeholder_id")
    assert hasattr(mapping, "action")
