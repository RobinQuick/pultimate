from schemas.template_spec import RgbColor, TemplateSpec, ThemeColors, ThemeFonts
from services.correction.engine import restyle_engine
from services.rules.base import FindingSpec
from pptx import Presentation


def create_bad_pptx(path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Bad Font Title"

    p = title.text_frame.paragraphs[0]
    p.font.name = "Comic Sans MS"  # Bad
    p.font.color.rgb = None  # Default

    # Save shape ID
    sid = title.shape_id

    prs.save(path)
    return str(sid)


def test_restyle_engine(tmp_path):
    # 1. Setup
    pptx_path = tmp_path / "bad.pptx"
    out_path = tmp_path / "fixed.pptx"
    shape_id = create_bad_pptx(pptx_path)

    template = TemplateSpec(
        name="Good Template",
        theme_fonts=ThemeFonts(major="Arial", minor="Inter"),
        theme_colors=ThemeColors(dark1=RgbColor(r=0, g=0, b=0)),
        masters=[],
    )

    findings = [
        FindingSpec(rule_id="FONT_MISMATCH", slide_index=0, element_id=shape_id, severity="HIGH", message="Bad Font")
    ]

    # 2. Run
    results = restyle_engine.apply_fixes(str(pptx_path), str(out_path), findings, template, mode="SAFE")

    # 3. Verify Result Object
    assert len(results) == 1
    assert results[0].status == "SUCCESS"
    assert "Changed font" in results[0].action_taken

    # 4. Verify PPTX Output
    prs = Presentation(str(out_path))
    slide = prs.slides[0]
    title = slide.shapes[0]  # Finding shape 0
    font_name = title.text_frame.paragraphs[0].font.name
    # FontFixer logic: uses 'minor' (Inter) as default fallback or 'major' depending on heuristic.
    # Our implementation used 'minor' (Inter).
    assert font_name == "Inter"
